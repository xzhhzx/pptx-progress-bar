# from collections.abc import Container  # (workaround for python 3.10)
import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

""" Default params """
# Tag name for progress bar, which is convenient for searching shapes on a slide by name
PROGRESS_BAR_TAG = "progress_bar_tag"
# The fixed name for the chapter slide of Microsoft Powerpoint (in Chinese).
CHAPTER_SLIDE_LAYOUT_NAME = "节标题"

class ColorsManager(object):
    """ This class is for managing and retrieving colors. Each chapter would be
        assigned with a different color, in a round-robin fashion.
    """
    def __init__(self, colors):
        self.ptr = 0
        self.colors = colors    # a set of available colors (in hex format) for different chapters
        self.rgb_colors = [RGBColor(*self._convert_hex_to_rgb(hex_color)) for hex_color in self.colors]

    def _convert_hex_to_rgb(self, hex_color):
        return (
            int(hex_color[:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:], 16)
        )

    def getCurrentColor(self):
        return self.rgb_colors[self.ptr]

    def changeToNextColor(self):
        self.ptr = (self.ptr + 1) % len(self.colors)

    def resetColor(self):
        self.ptr = 0


class ProgressBarTemplate(object):
    """ A class that represents the progress bar template that can be drawn on
        a PPT slide. To use this, first instantiate with the builder. Then use
        the method *drawAllBars()*/*removeAllBars()* to draw/remove progress
        bars on every page.
    """
    def __init__(self, builder):
        # Progress bar properties
        self.position = builder.position
        self.thk = builder.thk
        self.chapterColorsManager = builder.chapterColorsManager
        self.chapterColorsManagerBg = builder.chapterColorsManagerBg
        self.unit_size = builder.unit_size
        self.thk_bg = builder.thk_bg
        self.bg_margin = builder.bg_margin
        self.chapter_tuple_list = builder.chapter_tuple_list
        self.chapter_start_pages = [i[0] for i in self.chapter_tuple_list]
        self.num_pages_of_chapters = [i[0] - i[1] for i in zip(self.chapter_start_pages[1:], self.chapter_start_pages[:-1])]

        # Presentation properties
        self.prs = builder.prs
        self.W = builder.W
        self.H = builder.H


    def _appendRect(self, shapes, offset, delta):
        """ Append a foreground rectangle to the progress bar """
        if self.position in ['bottom', 'top']:
            rect = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                offset,
                (self.H - self.thk) if self.position == 'bottom' else 0,
                delta,
                self.thk
            )
        elif self.position in ['right', 'left']:
            rect = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                (self.W - self.thk) if self.position == 'right' else 0,
                offset,
                self.thk,
                delta
            )
        self._fillPureColor(rect, self.chapterColorsManager)
        return rect

    def _appendRectBg(self, shapes, offset, delta):
        """ Append a background rectangle to the progress bar
            (although this is a duplicate of _appendRect, I think in this way
            the code logic looks cleaner)
        """
        if self.position in ['bottom', 'top']:
            rect = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                offset,
                (self.H - self.thk_bg - self.bg_margin) if self.position == 'bottom' else self.bg_margin,
                delta,
                self.thk_bg
            )
        elif self.position in ['right', 'left']:
            rect = shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                (self.W - self.thk_bg - self.bg_margin) if self.position == 'right' else self.bg_margin,
                offset,
                self.thk_bg,
                delta
            )
        self._fillPureColor(rect, self.chapterColorsManagerBg)
        return rect

    def _fillPureColor(self, shape, colorFactory):
        """ Helper function: fill pure color to the shape with a ColorFactory """
        shape.fill.solid()
        shape.fill.fore_color.rgb = colorFactory.getCurrentColor()
        shape.line.fill.background()     # set line to transparent
        colorFactory.changeToNextColor()

    def _drawBarOnPage(self, page):
        """ Draw a progress bar for a certain page on slide """
        slide = self.prs.slides[page]
        group_shape = slide.shapes.add_group_shape([])
        group_shape.name = PROGRESS_BAR_TAG  # tag a name for the shape
        ptr = 0
        offset = 0
        page += 1 # TODO: this is ugly. Try fixing the offset issue with an alternative method

        # 1.Previous full chapters
        while self.chapter_tuple_list[ptr+1][0] < page:
            delta = self.unit_size * self.num_pages_of_chapters[ptr]
            self._appendRect(group_shape.shapes, offset, delta)
            offset += delta
            ptr += 1

        # 2.Current chapter
        delta_pages = page - self.chapter_tuple_list[ptr][0]
        delta = self.unit_size * delta_pages
        self._appendRect(group_shape.shapes, offset, delta)
        offset += delta

        # 3.Remaining all pages
        total_pages = self.chapter_start_pages[-1]
        delta = self.unit_size * (total_pages - page)
        self._appendRectBg(group_shape.shapes, offset, delta)
        offset += delta

        self.chapterColorsManager.resetColor()
        self.chapterColorsManagerBg.resetColor()

    def drawAllBars(self):
        """ Draw all progress bars for the whole presentation """
        for i in range(len(self.prs.slides)):
            self._drawBarOnPage(i)

    def removeAllBars(self):
        """ Remove all progress bars for the whole presentation """
        for slide in self.prs.slides:
            # Search shape named with "progress_bar" and remove
            for shape in slide.shapes:
                if shape.name.startswith(PROGRESS_BAR_TAG):
                    slide.shapes.element.remove(shape.element)

    class ProgressBarTemplateBuilder(object):
        """ Builder pattern"""
        def __init__(self, presentation):
            # Required params
            self.prs = presentation
            self.W = presentation.slide_width
            self.H = presentation.slide_height
            self.chapter_tuple_list = self._calculateChapterSegments()

            print('============= AUTO-DETECTED CHAPTERS =============')
            for i in self.chapter_tuple_list:
                print(i)

            # Optional params
            self.position = 'bottom'
            self.thk = Inches(0.3)
            self.bg_thickness_ratio = 0.5
            self.chapterColorsManager = ColorsManager(["540d6e", "ee4266", "ffd23f", "3bceac"])
            self.chapterColorsManagerBg = ColorsManager(["D8E1E9"])

        def _calculateChapterSegments(self):
            """ Pre-calculate all chapter segments. The first and last element
                represents the start/end of the whole presentation.
                FORMAT: each tuple represents (<chapter_page_index>, <chapter_name>)
            """
            chapter_tuple_list = [(0, "start")] # (start page)
            for idx, slide in enumerate(self.prs.slides):
                if slide.slide_layout.name == CHAPTER_SLIDE_LAYOUT_NAME:
                    chapter_tuple_list.append((idx, slide.shapes[0].text))
            chapter_tuple_list.append((len(self.prs.slides), "end")) # (end_page)
            return chapter_tuple_list

        def _checkNone(self, attr):
            return attr is None

        # TODOs: add validation check
        def setPosition(self, position):
            # reject setting if input is none (TODO: is this a good design?)
            if self._checkNone(position):
                return self
            if position not in ('top', 'bottom', 'left', 'right'):
                raise Exception("Input position should be one of: ('top', 'bottom', 'left', 'right')")
            self.position = position
            return self

        def setThickness(self, thk):
            self.thk = Inches(thk)
            return self

        def setBgThicknessRatio(self, bg_thickness_ratio):
            """ The background thickness ratio, relative to self.thk"""
            self.bg_thickness_ratio = bg_thickness_ratio
            return self

        def setColors(self, colors):
            self.chapterColorsManager = ColorsManager(colors)
            return self

        def setBgColor(self, bg_color):
            self.chapterColorsManagerBg = ColorsManager([bg_color])
            return self

        def build(self):
            num_pages = self.chapter_tuple_list[-1][0]
            if self.position in ['bottom', 'top']:
                self.unit_size = self.W / num_pages  # the base unit of the bar on one page
            elif self.position in ['right', 'left']:
                self.unit_size = self.H / num_pages
            self.thk_bg = self.bg_thickness_ratio * self.thk
            self.bg_margin = (self.thk - self.thk_bg) / 2   # the margin between foreground and background
            return ProgressBarTemplate(self)
